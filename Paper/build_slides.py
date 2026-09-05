# -*- coding: utf-8 -*-
"""Paper/build_slides.py — 논문 발표용 슬라이드(.pptx) 생성 (국문/영문).

16:9. 그림은 Paper/figures/ 의 논문용 그림을 그대로 재사용한다.
실행: C:/ml_env/Scripts/python Paper/build_slides.py
"""
import os
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Cm, Pt

sys.stdout.reconfigure(encoding="utf-8")
HERE = os.path.dirname(os.path.abspath(__file__))
FIG = os.path.join(HERE, "figures")

KO_FONT, EN_FONT = "맑은 고딕", "Segoe UI"
MAROON = RGBColor(0x8C, 0x1D, 0x40)
BEIGE = RGBColor(0xC8, 0xB4, 0x9A)
DARK = RGBColor(0x20, 0x20, 0x20)
GREY = RGBColor(0x60, 0x60, 0x60)
LIGHT = RGBColor(0xF2, 0xF2, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PINK = RGBColor(0xFD, 0xF0, 0xF0)
PINK2 = RGBColor(0xFD, 0xF2, 0xF5)
BLUEBG = RGBColor(0xEE, 0xF4, 0xFF)
HEADBG = RGBColor(0x44, 0x44, 0x44)
C_L = RGBColor(0xFF, 0x8C, 0x00)   # 인덕턴스 L — 주황
C_V = RGBColor(0x1F, 0x77, 0xB4)   # TENG V — 파랑
C_F = RGBColor(0xD6, 0x2D, 0x2D)   # 접촉력 — 빨강

SW, SH = Cm(33.867), Cm(19.05)     # 16:9
MARGIN = Cm(1.6)
BODY_TOP = Cm(3.3)
FULLW = SW - 2 * MARGIN

FONT = KO_FONT   # build() 에서 언어별로 교체


def _font(run, size=14, bold=False, color=DARK, italic=False):
    f = run.font
    f.name = FONT
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", FONT)


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Cm(0)
    tf.margin_top = tf.margin_bottom = Cm(0)
    return tf


def put(tf, text, size=14, bold=False, color=DARK, space_after=6,
        align=PP_ALIGN.LEFT, bullet=None, first=False, italic=False, line=1.25):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.line_spacing = line
    if bullet:
        r = p.add_run()
        r.text = bullet + "  "
        _font(r, size, True, MAROON)
    r = p.add_run()
    r.text = text
    _font(r, size, bold, color, italic)
    return p


def rect(slide, x, y, w, h, fill=LIGHT, radius=True):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    if radius:
        try:
            shp.adjustments[0] = 0.06
        except Exception:
            pass
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Cm(0.3)
    tf.margin_top = tf.margin_bottom = Cm(0.22)
    return shp, tf


def line(slide, x, y, w, color=MAROON, weight=2.5):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Pt(weight))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def table(slide, hdr, rows, cw, x0, y0, hdr_h=Cm(1.0), row_h=Cm(1.4),
          hdr_size=13, row_size=13, left_cols=(0,), accent_col=None,
          accent_row=None, gap=Cm(0.14)):
    """헤더 + 본문 행으로 구성된 간단한 표. accent_col/accent_row 는 강조 색을 입힌다."""
    x = x0
    for i, h in enumerate(hdr):
        fill = MAROON if accent_col == i else HEADBG
        shp, tf = rect(slide, x, y0, cw[i] - gap, hdr_h, fill)
        put(tf, h, hdr_size, True, WHITE, first=True, align=PP_ALIGN.CENTER,
            space_after=0, line=1.05)
        x += cw[i]
    y = y0 + hdr_h + Cm(0.15)
    for k, row in enumerate(rows):
        x = x0
        for i, txt in enumerate(row):
            hot = (accent_col == i) or (accent_row == k)
            shp, tf = rect(slide, x, y, cw[i] - gap, row_h, PINK2 if hot else LIGHT)
            put(tf, txt, row_size, hot, MAROON if hot else DARK, first=True,
                align=PP_ALIGN.LEFT if i in left_cols else PP_ALIGN.CENTER,
                space_after=0, line=1.2)
            x += cw[i]
        y += row_h + Cm(0.14)
    return y


class Deck:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width, self.prs.slide_height = SW, SH
        self.blank = self.prs.slide_layouts[6]
        self.n = 0

    def slide(self, title=None, sub=None):
        s = self.prs.slides.add_slide(self.blank)
        if title is not None:
            self.n += 1
            tf = textbox(s, MARGIN, Cm(1.0), SW - 2 * MARGIN - Cm(3), Cm(1.2))
            put(tf, title, 25, True, DARK, first=True, space_after=0)
            if sub:
                tf2 = textbox(s, MARGIN, Cm(2.22), SW - 2 * MARGIN - Cm(3), Cm(0.8))
                put(tf2, sub, 12.5, False, GREY, first=True, space_after=0)
            line(s, MARGIN, Cm(2.98), SW - 2 * MARGIN - Cm(4.2))
            line(s, SW - MARGIN - Cm(3.6), Cm(2.98), Cm(3.6), BEIGE, 2.5)
            tfn = textbox(s, SW - MARGIN - Cm(2.4), SH - Cm(1.15), Cm(2.4), Cm(0.7))
            put(tfn, str(self.n), 11, False, GREY, first=True,
                align=PP_ALIGN.RIGHT, space_after=0)
        return s

    def pic(self, slide, name, x, y, w=None, h=None):
        p = os.path.join(FIG, name)
        if w is not None:
            return slide.shapes.add_picture(p, x, y, width=w)
        return slide.shapes.add_picture(p, x, y, height=h)


def build(L):
    global FONT
    FONT = KO_FONT if L == "ko" else EN_FONT

    def t(ko, en):
        return ko if L == "ko" else en

    d = Deck()

    # ── 1. 표지 ─────────────────────────────────────────────────────────────
    s = d.slide()
    rect(s, Cm(0), Cm(0), SW, Cm(0.45), MAROON, radius=False)
    tf = textbox(s, MARGIN, Cm(4.8), FULLW, Cm(5.2))
    put(tf, t("단일 전극 소프트 멀티모달 센서의 시분할 측정과",
              "Time-Division Measurement and Embedded-AI Signal"),
        t(32, 30), True, DARK, first=True, space_after=6)
    put(tf, t("임베디드 AI 기반 실시간 신호 디커플링",
              "Decoupling for a Soft Single-Electrode Multimodal Sensor"),
        t(32, 30), True, DARK, space_after=14)
    put(tf, t("인장 · 근접 · 접촉력의 동시 추정",
              "Simultaneous Estimation of Strain, Proximity, and Contact Force"),
        18, False, MAROON, space_after=0)
    line(s, MARGIN, Cm(11.4), Cm(8))
    tf = textbox(s, MARGIN, Cm(12.1), FULLW, Cm(3))
    put(tf, t("문희준", "Heejun Moon"), 20, True, DARK, first=True, space_after=6)
    put(tf, t("고려대학교 기계공학과  BioRobotics & Control Lab (BiRC)",
              "BioRobotics and Control Lab (BiRC),  Dept. of Mechanical Engineering,  "
              "Korea University"),
        14, False, GREY, space_after=0)

    # ── 2. 목차 ─────────────────────────────────────────────────────────────
    s = d.slide(t("목차", "Outline"))
    items = [
        ("1", t("연구 배경", "Motivation"),
         t("멀티모달 인지의 필요성 · 각 모달리티의 역할",
           "Why multimodal perception, and what each modality is for")),
        ("2", t("문제 정의", "Problem Statement"),
         t("신호 결합과 기존 연구의 한계",
           "Signal coupling and the limitations of prior work")),
        ("3", t("센서 설계", "Sensor Design"),
         t("단일 전극 EGaIn 나선 코일 · 제작 공정",
           "Single-electrode EGaIn spiral coil and fabrication")),
        ("4", t("측정 시스템", "Measurement System"),
         t("1 ms TDM 아키텍처 · PCB · 시험 플랫폼",
           "1 ms TDM architecture, PCB and test platform")),
        ("5", t("신호 특성화", "Signal Characterization"),
         t("(ε, d) 응답 곡면 · 접촉 영역 점탄성",
           "(strain, d) response surfaces and contact-regime viscoelasticity")),
        ("6", t("디커플링 모델", "Decoupling Models"),
         t("2단계 구조 → SLS-EMA + Gate-MoE",
           "Two-stage structure → SLS-EMA + Gate-MoE")),
        ("7", t("임베디드 구현 및 실시간 검증",
                "Embedded Implementation and Real-Time Validation"),
         t("STM32 온보드 추론 260 µs", "260 µs on-board inference on the STM32")),
        ("8", t("결론 및 향후 계획", "Conclusion and Future Work"), ""),
    ]
    y = BODY_TOP + Cm(0.4)
    for num, t1, t2 in items:
        shp, tf = rect(s, MARGIN, y, Cm(1.5), Cm(1.28), MAROON)
        put(tf, num, 16, True, WHITE, first=True, align=PP_ALIGN.CENTER, space_after=0)
        tf = textbox(s, MARGIN + Cm(2.0), y + Cm(0.06), Cm(28), Cm(1.28))
        put(tf, t1, 17, True, DARK, first=True, space_after=1)
        if t2:
            put(tf, t2, 12, False, GREY, space_after=0)
        y += Cm(1.72)

    # ── 3. 연구 배경 ────────────────────────────────────────────────────────
    s = d.slide(t("연구 배경", "Motivation"),
                t("로봇 표면은 주변 상황과 자기 자신의 상태를 함께 인지해야 한다",
                  "A robot surface must perceive both its surroundings and its own state"))
    boxes = [
        (t("접촉 이전", "Before contact"),
         t("물체가 다가오는 것을\n미리 감지해야 한다",
           "The approach of an object\nmust be sensed in advance"),
         t("근접 (Proximity)", "Proximity"), C_V),
        (t("접촉 순간", "At contact"),
         t("얼마나 세게 닿아 있는지\n알아야 한다",
           "How hard the contact is\nmust be known"),
         t("접촉력 (Force)", "Contact force"), C_F),
        (t("상시", "At all times"),
         t("센서가 붙은 몸체 자체가\n변형한다",
           "The body carrying the sensor\nitself deforms"),
         t("인장 (Strain)", "Strain"), C_L),
    ]
    w = (FULLW - Cm(1.2)) / 3
    for i, (when, why, what, col) in enumerate(boxes):
        x = MARGIN + i * (w + Cm(0.6))
        shp, tf = rect(s, x, BODY_TOP + Cm(0.5), w, Cm(5.6), LIGHT)
        put(tf, when, 13, True, col, first=True, space_after=8)
        put(tf, why, 15, False, DARK, space_after=12, line=1.35)
        put(tf, what, 17, True, col, space_after=0)
    tf = textbox(s, MARGIN, BODY_TOP + Cm(7.0), FULLW, Cm(4))
    put(tf, t("그런데 실제 상호작용에서는 세 자극이 동시에 작용한다.",
              "In real interaction, however, all three stimuli act simultaneously."),
        18, True, DARK, first=True, space_after=10)
    put(tf, t("→ 단일 물리량만 측정하는 센서는 현실 상황에서 활용도가 제한된다.",
              "→ A sensor measuring only one physical quantity has limited practical "
              "utility."), 15, False, GREY, space_after=6)
    put(tf, t("→ 서로 다른 자극이 하나의 전기적 응답에 겹쳐 들어온다  =  신호 결합 (signal coupling)",
              "→ Distinct stimuli superimpose onto a single electrical response  =  "
              "signal coupling"), 15, True, MAROON, space_after=0)

    # ── 4. 각 모달리티의 필요성 ─────────────────────────────────────────────
    s = d.slide(t("각 모달리티는 왜 필요한가", "Why Each Modality Is Needed"),
                t("선행 연구가 말하는 근거", "The case made by prior work"))
    rows = [
        (t("근접", "Proximity"), C_V,
         t("충돌이 일어나기 전에 개입할 수 있는 유일한 수단",
           "The only means of intervening before a collision occurs"),
         t("• 안전한 pHRI: 접촉 후 반응으로는 불충분 → 접근을 예측하고 궤적 수정 [6],[7]\n"
           "• 로봇 팔 외피 전면 분포 배치로 표면 전체 충돌 회피 [8], 준전신 감지 [9]\n"
           "• 파지 시 손이 닿기 전 자세 정렬",
           "• Safe pHRI: reacting after contact is insufficient → anticipate the approach "
           "and correct the trajectory [6], [7]\n"
           "• Sensors over the whole arm exterior for full-surface collision avoidance "
           "[8]; quasi whole-body sensing [9]\n"
           "• Pre-shaping the hand before it reaches the object")),
        (t("접촉력", "Contact force"), C_F,
         t("파지 안정성은 가하는 힘의 크기에 직접 좌우",
           "Grasp stability depends directly on the applied force"),
         t("• 촉각 피드백이 없으면 물체를 놓치거나 과도한 힘으로 손상 [10]\n"
           "• 미끄러짐은 빠르게 검출·보정하지 않으면 곧바로 파지 실패 [11]\n"
           "• 근접이 “언제 닿는가”라면, 접촉력은 “얼마나 세게 닿았는가” [12]",
           "• Without tactile feedback an object is dropped, or damaged by excessive "
           "force [10]\n"
           "• Slip leads immediately to grasp failure unless detected and corrected "
           "quickly [11]\n"
           "• Proximity answers when contact occurs; force answers how hard [12]")),
        (t("인장", "Strain"), C_L,
         t("센서가 붙은 몸체가 변형하는 경우 필수",
           "Essential whenever the body carrying the sensor deforms"),
         t("• 소프트 로봇은 사실상 무한 자유도 → 개루프 제어로 정확한 작업 곤란 [13]\n"
           "• 자기수용감각이 제어 루프를 닫는 전제 조건 [14]\n"
           "• 액체금속 변형 센서 + 소프트 크롤링 로봇 폐루프 제어 사례 [15]",
           "• Soft robots have effectively infinite DOF → accurate open-loop control is "
           "very difficult [13]\n"
           "• Proprioception is a precondition for closing the control loop [14]\n"
           "• Liquid-metal strain sensors enabling closed-loop soft crawling robots [15]")),
    ]
    y = BODY_TOP + Cm(0.3)
    for name, col, head, body in rows:
        shp, tf = rect(s, MARGIN, y, Cm(3.6), Cm(3.5), col)
        put(tf, name, t(20, 17), True, WHITE, first=True, align=PP_ALIGN.CENTER,
            space_after=0, line=1.1)
        tf = textbox(s, MARGIN + Cm(4.1), y + Cm(0.15), FULLW - Cm(4.4), Cm(3.5))
        put(tf, head, 15, True, DARK, first=True, space_after=7)
        put(tf, body, t(12.5, 11.5), False, GREY, space_after=0, line=1.35)
        y += Cm(3.85)

    # ── 5. 왜 함께 측정해야 하는가 ──────────────────────────────────────────
    s = d.slide(t("왜 한 지점에서 함께 측정해야 하는가",
                  "Why They Must Be Measured Together, at One Point"),
                t("세 모달리티는 독립적이지 않다",
                  "The three modalities are not independent"))
    tf = textbox(s, MARGIN, BODY_TOP + Cm(0.4), FULLW, Cm(2))
    put(tf, t("변형하는 표면 위에서는 근접·접촉 신호의 기준 자체가 변형 상태에 따라 달라진다.",
              "On a deforming surface, the very reference for the proximity and contact "
              "signals shifts with the deformation state."),
        17, True, DARK, first=True, space_after=0)
    shp, tf = rect(s, MARGIN, BODY_TOP + Cm(1.9), FULLW, Cm(3.5), PINK)
    put(tf, t("실측: 인장 상태가 압력 민감도를 6.6배까지 바꾼다  (교차 민감도)",
              "Measured: the strain state changes the pressure sensitivity by up to "
              "6.6×  (cross-sensitivity)"),
        17, True, MAROON, first=True, space_after=10)
    put(tf, t("ε = 0 %  에서 10 N 가압 → ΔR/R₀ 변화 폭  3.95 %p",
              "At ε = 0 %,  loading to 10 N   →  ΔR/R₀ span  3.95 %p"),
        15, False, DARK, space_after=4)
    put(tf, t("ε = 30 % 에서 9.2 N 가압 → ΔR/R₀ 변화 폭  26.14 %p",
              "At ε = 30 %, loading to 9.2 N  →  ΔR/R₀ span  26.14 %p"),
        15, False, DARK, space_after=0)
    tf = textbox(s, MARGIN, BODY_TOP + Cm(5.9), FULLW, Cm(5))
    put(tf, t("→ 변형률을 모르면, 같은 저항 변화가 어느 정도의 접촉력인지 확정할 수 없다.",
              "→ Without knowing the strain, a given resistance change cannot be mapped "
              "to a contact force."),
        16, True, DARK, first=True, space_after=12)
    put(tf, t("세 물리량을 각각 별개의 센서로 측정하더라도 이 결합은 사라지지 않는다.",
              "Measuring the three quantities with three separate sensors does not remove "
              "this coupling."),
        15, False, GREY, space_after=5, bullet="•")
    put(tf, t("오히려 서로 다른 위치에서 측정된 값을 정합해야 하는 문제가 추가된다.",
              "It merely adds the problem of registering values measured at different "
              "locations."),
        15, False, GREY, space_after=12, bullet="•")
    put(tf, t("결합이 발생하는 바로 그 지점에서 세 신호를 동시에 취득하고 함께 디커플링한다.",
              "Acquire the three signals at the very point where the coupling arises, and "
              "decouple them jointly."),
        16, True, MAROON, space_after=0, bullet="▶")

    # ── 6. 기존 연구의 한계 ─────────────────────────────────────────────────
    s = d.slide(t("기존 연구의 한계와 본 연구의 대응",
                  "Limitations of Prior Work and How This Work Responds"))
    hdr = [t("선행 접근", "Prior approach"), t("한계", "Limitation"),
           t("본 연구의 대응", "Response in this work")]
    data = [
        (t("강체 근접 센서 분포 배치", "Distributed rigid proximity sensors"),
         t("곡면·대변형 표면에 밀착 불가\n배선 수 ∝ 센서 수",
           "Cannot conform to curved or deformable\nsurfaces; wire count ∝ sensor count"),
         t("두 가닥 단일 전극의\n평면·신축 코일 하나",
           "A single planar, stretchable coil\nwith a two-wire electrode")),
        (t("정전용량형 근접+촉각 센서", "Capacitive proximity + tactile sensors"),
         t("인장과 근접이 같은 정전용량 변화로\n나타나 원리적 분리 곤란",
           "Strain and proximity appear as the same\ncapacitance change — hard to separate"),
         t("L은 (ε, d), R은 ε에만 의존하는\n물리적 비대칭 이용",
           "Exploits the asymmetry: L depends on\n(ε, d) while R depends on ε alone")),
        (t("인덕티브 소프트 센서 (단일 모달)", "Inductive soft sensors (single modality)"),
         t("디커플링 난이도로\n인장 또는 근접 중 하나만 사용",
           "Use only strain or only proximity\nbecause decoupling is hard"),
         t("동일 코일에서 3종 동시 추정",
           "Estimates all three simultaneously\nfrom the same coil")),
        (t("다중 소자 하이브리드 통합", "Hybrid integration of multiple devices"),
         t("소자 수만큼 배선·두께·크로스토크 증가",
           "Wiring, thickness and cross-talk\ngrow with the device count"),
         t("소자를 늘리지 않고\n1 ms TDM으로 시분할 사용",
           "Time-shares one device with a 1 ms\nTDM schedule instead of adding devices")),
        (t("LCR 미터 기반 인덕턴스 계측", "Inductance measured with an LCR meter"),
         t("탁상형 계측기 필요\n이동형·착용형 응용 불가",
           "Requires benchtop instrumentation;\nnot mobile or wearable"),
         t("LDC1614를 실은 65 × 52 mm 보드",
           "Replaced by a 65 × 52 mm board\ncarrying an LDC1614")),
        (t("PC 기반 AI 디커플링", "PC-based AI decoupling"),
         t("USB 왕복 20–40 ms, 지터 ±5–15 ms\n→ 되먹임 제어 불가",
           "20–40 ms USB round trip, ±5–15 ms jitter\n→ feedback control impossible"),
         t("MCU 온보드 추론 260 µs\n지연 산포 약 2 µs",
           "On-board MCU inference at 260 µs\nwith ~2 µs spread")),
    ]
    cw = [Cm(8.0), Cm(11.5), Cm(11.1)]
    table(s, hdr, data, cw, MARGIN, BODY_TOP + Cm(0.2),
          hdr_h=Cm(1.0), row_h=Cm(1.72), hdr_size=13,
          row_size=t(11.5, 11), left_cols=(0, 1, 2), accent_col=2)

    # ── 7. 접근 요약 ────────────────────────────────────────────────────────
    s = d.slide(t("본 연구의 접근", "Our Approach"),
                t("세 가지 설계 결정", "Three design decisions"))
    cards = [
        (t("① 단일 전극", "① Single electrode"),
         t("소자를 늘리는 대신\n하나의 소자를 시간축에서 나눠 쓴다",
           "Divide one device in time\nrather than multiplying devices"),
         t("• EGaIn 나선 코일 1개, 배선 2가닥\n• 소자 간 크로스토크 원리적으로 없음\n"
           "• 배선 수는 모달리티 수와 무관하게 2\n• 평면성·신축성 유지",
           "• One EGaIn spiral coil, two wires\n• Device cross-talk cannot arise\n"
           "• Wire count stays 2 regardless of the\n   number of modalities\n"
           "• Planarity and stretchability preserved")),
        (t("② 1 ms TDM", "② 1 ms TDM"),
         t("AC 여기 · DC 여기 · 무여기를\n순차 인가",
           "Apply AC excitation, DC excitation\nand no excitation in sequence"),
         t("• AC → 공진 주파수 (인덕턴스 L)\n• DC → 전압 분배 (저항 R)\n"
           "• 무여기 → 부유 상태 (TENG 전압 V)\n• 1 kHz 취득, CPU 유휴율 93.9 %",
           "• AC → resonant frequency (L)\n• DC → voltage divider (R)\n"
           "• None → floating coil (TENG V)\n• 1 kHz acquisition, 93.9 % CPU idle")),
        (t("③ 임베디드 AI", "③ Embedded AI"),
         t("계측과 추론을 모두\nMCU 안에서 끝낸다",
           "Complete both measurement and\ninference inside the MCU"),
         t("• 해석적 역산은 특이점에서 발산\n• PC 추론은 USB 왕복 20–40 ms\n"
           "• 1,413 파라미터 모델을 STM32에 임베딩\n• 실측 260 µs, 1 ms 주기 내 완결",
           "• Analytical inversion diverges at\n   singularities\n"
           "• PC inference costs a 20–40 ms round trip\n"
           "• 1,413-parameter model on the STM32\n"
           "• Measured 260 µs, inside the 1 ms cycle")),
    ]
    w = (FULLW - Cm(1.2)) / 3
    for i, (t1, t2, t3) in enumerate(cards):
        x = MARGIN + i * (w + Cm(0.6))
        shp, tf = rect(s, x, BODY_TOP + Cm(0.4), w, Cm(1.3), MAROON)
        put(tf, t1, 18, True, WHITE, first=True, align=PP_ALIGN.CENTER, space_after=0)
        shp, tf = rect(s, x, BODY_TOP + Cm(1.9), w, Cm(2.4), LIGHT)
        put(tf, t2, 14, True, DARK, first=True, space_after=0, line=1.3)
        tf = textbox(s, x + Cm(0.3), BODY_TOP + Cm(4.7), w - Cm(0.6), Cm(6))
        put(tf, t3, t(12.5, 12), False, GREY, first=True, space_after=0, line=1.5)

    # ── 8. 센서 설계 ────────────────────────────────────────────────────────
    s = d.slide(t("센서 설계와 제작", "Sensor Design and Fabrication"),
                t("EGaIn 액체금속 평면 직사각형 나선 코일 (초기 길이 120 mm, 인장 0–30 %)",
                  "EGaIn liquid-metal planar rectangular spiral coil "
                  "(120 mm initial length, 0–30 % strain)"))
    d.pic(s, "fig0_sensor.png", MARGIN, BODY_TOP + Cm(0.2), w=Cm(20.5))
    tf = textbox(s, MARGIN + Cm(21.2), BODY_TOP + Cm(0.5), Cm(9.6), Cm(12))
    put(tf, t("핵심 공정: 브리징", "Key step: bridging"), 16, True, MAROON,
        first=True, space_after=8)
    put(tf, t("나선의 안쪽 끝 단자를 권선 위로 넘겨 바깥으로 빼낸다.",
              "The inner terminal of the spiral is routed out over the windings."),
        13, False, DARK, space_after=8, line=1.4)
    put(tf, t("이 단계가 없으면 내부 단자에 접근할 수 없어 관통 배선이나 대향 전극이 필요해지고, "
              "그 순간 평면성과 신축성이 훼손된다.",
              "Without it the inner terminal cannot be reached without a "
              "through-connection or a counter electrode — which would immediately "
              "compromise planarity and stretchability."),
        t(12.5, 12), False, GREY, space_after=16, line=1.4)
    put(tf, t("압력 민감도 향상: 돌기 구조", "Raising pressure sensitivity: bump array"),
        16, True, MAROON, space_after=8)
    put(tf, t("직경 5 mm · 높이 0.5 mm, 필름 마스킹 + 어플리케이터",
              "5 mm diameter, 0.5 mm height, formed by film masking and an applicator"),
        t(12.5, 12), False, GREY, space_after=8, line=1.4)
    put(tf, t("ΔR/R₀ 100 % 도달 하중\n210 N  →  20 N  (약 10배)",
              "Load for a 100 % ΔR/R₀ response\n210 N  →  20 N  (about 10×)"),
        14, True, DARK, space_after=0, line=1.4)

    # ── 9. 응답 시그니처 ────────────────────────────────────────────────────
    s = d.slide(t("전자기적 특성 — 자극별 응답 시그니처",
                  "Electromagnetic Characteristics — Response Signature"))
    hdr = [t("자극", "Stimulus"), "L", "R_s (AC)", "R_DC", "V_TENG",
           t("지배 기구", "Dominant mechanism")]
    data = [
        (t("인장", "Stretching"), "↑", "↑", "↑", "—",
         t("길이 ↑, 단면적 ↓ → R ↑ ;  코일 기하 변화 → L ↑",
           "length ↑, section ↓ → R ↑ ;  coil geometry → L ↑")),
        (t("도체 접근", "Conductor approach"), "↓", "↑", "—", "—",
         t("와전류 손실 → 반사 임피던스 (신호 감쇠)",
           "eddy-current loss → reflected impedance (attenuation)")),
        (t("유전체 접근", "Dielectric approach"), "↑", "—", "—", "—",
         t("유전율 변화 / 전기장 재분포 (신호 증폭)",
           "permittivity change / field redistribution (amplification)")),
        (t("대전체 접촉", "Contact (charged body)"), "—", "—", "—", "↑",
         t("마찰전기 전하 이동", "triboelectric charge transfer")),
        (t("압축 (돌기)", "Compression (bumps)"), "↓*", "↑", "↑", "—",
         t("국소 단면적 ↓ ;  Q ↓ → 겉보기 L ↑",
           "local section ↓ ;  Q ↓ → apparent L ↑")),
    ]
    cw = [Cm(5.0), Cm(2.2), Cm(3.0), Cm(2.6), Cm(3.0), Cm(14.8)]
    y = table(s, hdr, data, cw, MARGIN, BODY_TOP + Cm(0.5),
              hdr_h=Cm(1.0), row_h=Cm(1.4), hdr_size=13,
              row_size=t(13, 12), left_cols=(0, 5))
    tf = textbox(s, MARGIN, y + Cm(0.25), FULLW, Cm(2))
    put(tf, t("도체와 유전체에서 ΔL의 부호가 반대  →  단순한 방해 요인이 아니라, "
              "물체의 전기적 성질을 식별할 수 있는 정보",
              "ΔL has opposite signs for conductors and dielectrics  →  not a nuisance, "
              "but information for identifying an object's electrical property"),
        14, True, MAROON, first=True, space_after=0)

    # ── 10. TDM ─────────────────────────────────────────────────────────────
    s = d.slide(t("측정 시스템 — 1 ms 시분할 측정 (TDM)",
                  "Measurement System — 1 ms Time-Division Measurement"))
    d.pic(s, "fig1_tdm.png", Cm((33.867 - 17.2) / 2), BODY_TOP + Cm(0.1), h=Cm(11.8))
    tf = textbox(s, MARGIN, SH - Cm(2.9), FULLW, Cm(2))
    put(tf, t("TIM7 인터럽트로 두 경로 동시 시작 → I²C DMA(LDC1614 28비트, ~200 µs)와 "
              "ADC DMA(TENG 150 µs → R 150 µs) 병렬 실행",
              "A TIM7 interrupt starts both paths → I²C DMA (LDC1614, ~200 µs) in "
              "parallel with ADC DMA (TENG 150 µs → R 150 µs)"),
        13.5, False, GREY, first=True, space_after=5)
    put(tf, t("모든 전송이 DMA → 블로킹 대기 없음 · CPU 유휴율 93.9 % · 고정 샘플링 주기로 "
              "AI 결합 시 지터 방지",
              "All transfers use DMA → no blocking wait · 93.9 % CPU idle · a fixed "
              "period prevents jitter when coupled to the AI model"),
        13.5, False, GREY, space_after=0)

    # ── 11. 하드웨어 ────────────────────────────────────────────────────────
    s = d.slide(t("하드웨어 구성", "Hardware"),
                t("전용 PCB와 자체 제작 시험 플랫폼",
                  "A dedicated PCB and an in-house test platform"))
    groups = [
        (t("다신호 스위칭 PCB", "Multisignal switching PCB"), [
            t("STM32G473CBT6 — Cortex-M4F 170 MHz, Flash 128 KB, SRAM 32 KB",
              "STM32G473CBT6 — Cortex-M4F 170 MHz, 128 KB flash, 32 KB SRAM"),
            t("LDC1614 — 28비트, 1 kHz–10 MHz, I²C2 @ 0x2A",
              "LDC1614 — 28-bit, 1 kHz–10 MHz, I²C2 at 0x2A"),
            t("ADG734 — 4채널 아날로그 스위치, t_on 29 ns, BW 200 MHz",
              "ADG734 — 4-channel analog switch, t_on 29 ns, 200 MHz BW"),
            t("65 × 52 mm 6층 기판, 최대 4채널 센서 어레이",
              "65 × 52 mm six-layer board, up to a 4-channel array"),
            t("TENG 경로 TVS 다이오드 보호", "TVS diodes protecting the TENG path"),
        ]),
        (t("5축 시험 플랫폼", "Five-axis test platform"), [
            t("XA/XB/YA/YB 4축 대칭 인장 + 독립 Z축 근접",
              "Four symmetric tensile axes (XA/XB/YA/YB) + an independent Z proximity axis"),
            t("700 × 700 mm 프레임, Nema23 + TB6600, Arduino Uno",
              "700 × 700 mm frame, Nema23 + TB6600 drivers, Arduino Uno"),
            t("320 step/mm, PyQt5 GUI 시퀀스 구동",
              "320 steps/mm, PyQt5 GUI driving the sequence"),
            t("6축 F/T 센서(CAN)로 접촉력 참값 동시 기록",
              "Six-axis F/T sensor over CAN recording ground-truth force"),
            t("PC-MCU 시각 동기 + 가감속 구간 is_steady 자동 태깅",
              "PC-MCU time sync + automatic is_steady tagging of accel/decel transients"),
        ]),
    ]
    for col, (title, its) in enumerate(groups):
        x = MARGIN + col * (FULLW / 2 + Cm(0.4))
        w = FULLW / 2 - Cm(0.4)
        shp, tf = rect(s, x, BODY_TOP + Cm(0.4), w, Cm(1.15), MAROON)
        put(tf, title, 16, True, WHITE, first=True, align=PP_ALIGN.CENTER, space_after=0)
        tf = textbox(s, x + Cm(0.4), BODY_TOP + Cm(2.0), w - Cm(0.8), Cm(7))
        for k, it in enumerate(its):
            put(tf, it, t(13.5, 12.5), False, DARK, first=(k == 0), space_after=11,
                bullet="•", line=1.35)
    shp, tf = rect(s, MARGIN, BODY_TOP + Cm(9.6), FULLW, Cm(2.4), PINK)
    put(tf, t("왜 상용 인장 시험기를 쓰지 않았는가",
              "Why not a commercial tensile tester?"),
        14, True, MAROON, first=True, space_after=6)
    put(tf, t("한쪽 축만 구동하면 시편이 늘어나며 중심점이 이동한다. 근접 표적은 코일 중심 "
              "바로 위에 있어야 하므로, 중심 이동이 곧 근접 거리 라벨의 오차가 된다.",
              "In single-sided tension the specimen centre translates as it elongates. "
              "The proximity target must sit directly above the coil centre, so that "
              "translation becomes an error in the distance label."),
        13, False, DARK, space_after=0, line=1.3)

    # ── 12. 데이터셋 ────────────────────────────────────────────────────────
    s = d.slide(t("데이터셋", "Datasets"),
                t("모델 입력은 언제나 (ΔL/L₀, ΔR/R₀)뿐 — 위치·힘은 정답 라벨로만 사용",
                  "Model inputs are always only (ΔL/L₀, ΔR/R₀); positions and forces are "
                  "used solely as ground-truth labels"))
    hdr = ["", t("프로토콜", "Protocol"), t("규모", "Size"), t("용도", "Purpose")]
    data = [
        ("D1", t("변형 37단계 + 근접 50→0 mm 연속 (비동기)",
                 "37 strain levels + continuous 50→0 mm proximity (asynchronous)"),
         t("≈ 78,000행", "≈ 78,000 rows"),
         t("초기 2단 PINN 학습", "Initial two-stage PINN training")),
        ("D2", t("P1: 변형 19단계 × 근접 스윕 / P2: 근접 13단계 × 변형 스윕",
                 "P1: 19 strain levels × proximity sweep\nP2: 13 proximity levels × strain "
                 "sweep"),
         t("금속 45,633행\n손 37,116행", "metal 45,633 rows\nhand 37,116 rows"),
         t("(ε, d) 응답 곡면 특성화", "Characterizing the (ε, d) response surface")),
        ("D3", t("연속 동시 스윕 2세션, PC-MCU 동기",
                 "Two continuous simultaneous sweeps, PC-MCU synchronized"),
         t("188,232행\n→ 전처리 후 33,450", "188,232 rows\n→ 33,450 after preprocessing"),
         t("근접+인장 디커플러 학습", "Training the proximity + strain decoupler")),
        ("D4", t("변형 19단계 × 근접→접촉 왕복, F/T 동시 기록 (최대 10.1 N)",
                 "19 strain levels × proximity-to-contact traverse with F/T recording "
                 "(up to 10.1 N)"),
         t("51,858행 / 569 s", "51,858 rows / 569 s"),
         t("접촉 물리 분석 · MoE 학습", "Contact physics analysis · MoE training")),
        ("D5", t("실기 검증 3세션 (Part1 근접 / Part2 인장 / Part3 동시)",
                 "Three on-hardware sessions (Part 1 proximity / Part 2 strain / "
                 "Part 3 both)"),
         t("26.7k / 34.0k / 32.6k행", "26.7k / 34.0k / 32.6k rows"),
         t("임베디드 실시간 검증", "Embedded real-time validation")),
    ]
    cw = [Cm(2.4), Cm(13.6), Cm(7.2), Cm(7.4)]
    table(s, hdr, data, cw, MARGIN, BODY_TOP + Cm(0.4),
          hdr_h=Cm(0.95), row_h=Cm(1.85), hdr_size=13,
          row_size=t(12.5, 11.5), left_cols=(1, 3), accent_col=0)

    # ── 13. 핵심 근거 ───────────────────────────────────────────────────────
    s = d.slide(t("핵심 근거 — (ε, d) 응답 곡면",
                  "Key Evidence — the (ε, d) Response Surfaces"),
                t("디커플링이 가능한 물리적 이유",
                  "The physical reason decoupling is possible"))
    d.pic(s, "fig2_surfaces.png", MARGIN, BODY_TOP + Cm(0.2), w=Cm(20.0))
    tf = textbox(s, MARGIN + Cm(20.8), BODY_TOP + Cm(0.6), Cm(10.0), Cm(12))
    put(tf, t("금속 표적", "Metal target"), 16, True, MAROON, first=True, space_after=8)
    put(tf, t("ΔR 등고선이 근접거리 축과 거의 완전히 평행\n→ 저항은 인장에만 의존",
              "ΔR contours are essentially parallel to the distance axis\n"
              "→ resistance depends on strain alone"),
        t(13, 12), False, DARK, space_after=10, line=1.4)
    put(tf, t("ΔL 등고선은 d < 10 mm에서 급격히 휘어짐\n→ 같은 인장에서도 −15.7 % ~ +14.5 %로 "
              "부호까지 바뀜",
              "ΔL contours bend sharply for d < 10 mm\n"
              "→ at fixed strain it swings from −15.7 % to +14.5 %, changing sign"),
        t(13, 12), False, DARK, space_after=16, line=1.4)
    put(tf, t("손(유전체) 표적", "Hand (dielectric) target"), 16, True, MAROON,
        space_after=8)
    put(tf, t("두 신호 모두 등고선이 거의 수직\n→ 이 거리 범위에서 인덕턴스에 실린 근접 정보가 미약",
              "Contours of both signals are nearly vertical\n"
              "→ little proximity information is carried by the inductance here"),
        t(13, 12), False, DARK, space_after=16, line=1.4)
    shp, tf2 = rect(s, MARGIN + Cm(20.8), BODY_TOP + Cm(9.0), Cm(10.0), Cm(2.6), PINK)
    put(tf2, t("∇L 과 ∇R 이 전역적으로 선형 독립\n→ 디커플링 가능성의 직접적 근거",
               "∇L and ∇R are globally linearly independent\n"
               "→ direct evidence that decoupling is feasible"),
        t(13.5, 12.5), True, MAROON, first=True, align=PP_ALIGN.CENTER,
        space_after=0, line=1.35)

    # ── 14. 해석적 역산의 한계 ──────────────────────────────────────────────
    s = d.slide(t("그렇다면 해석적으로 역산하면 되지 않는가",
                  "Why Not Simply Invert the Relations Analytically?"),
                t("야코비안 역산이 실패하는 두 가지 이유",
                  "Two reasons Jacobian inversion fails"))
    shp, tf = rect(s, MARGIN, BODY_TOP + Cm(0.4), FULLW, Cm(2.3), LIGHT)
    put(tf, "[ΔR, ΔL]ᵀ  =  J · [Δε, Δd]ᵀ        →        "
            "[ε, d]ᵀ  =  [ε₀, d₀]ᵀ + J⁻¹ · [ΔR, ΔL]ᵀ",
        t(18, 16), True, DARK, first=True, align=PP_ALIGN.CENTER, space_after=0)
    items = [
        (t("① 절단 오차", "① Truncation error"),
         t("근거리에서는 M² 항이, 대변형에서는 sinh⁻¹ · ln 항이 지배하는 강한 비선형 영역\n"
           "→ 선형 근사가 국소적으로만 유효",
           "The near field is dominated by the M² term and large strain by sinh⁻¹ and ln "
           "terms\n→ the linear approximation is valid only locally")),
        (t("② 특이점 · 직교성 상실", "② Singularity / loss of orthogonality"),
         t("∇R 과 ∇L 의 방향이 가까워지면 det(J) → 0\n→ J⁻¹ 발산, 작은 측정 잡음이 크게 증폭",
           "As the directions of ∇R and ∇L converge, det(J) → 0\n"
           "→ J⁻¹ diverges and small measurement noise is greatly amplified")),
    ]
    y = BODY_TOP + Cm(3.2)
    for t1, t2 in items:
        shp, tf = rect(s, MARGIN, y, Cm(8.6), Cm(2.5), MAROON)
        put(tf, t1, t(16, 15), True, WHITE, first=True, align=PP_ALIGN.CENTER,
            space_after=0, line=1.15)
        tf = textbox(s, MARGIN + Cm(9.2), y + Cm(0.35), FULLW - Cm(9.5), Cm(2.5))
        put(tf, t2, t(14.5, 13.5), False, DARK, first=True, space_after=0, line=1.4)
        y += Cm(3.0)
    shp, tf = rect(s, MARGIN, y + Cm(0.3), FULLW, Cm(3.0), PINK)
    put(tf, t("실측: 교차항을 포함한 물리식 직접 역산 →  R² = 0.26   "
              "(품질 낮은 데이터 포함 시 ε ≈ 0 부근에서 R² = −75 까지 발산)",
              "Measured: direct inversion of the physics expression with the cross term "
              "→  R² = 0.26   (diverging to R² = −75 near ε ≈ 0 with lower-quality data)"),
        t(15, 14), True, MAROON, first=True, space_after=8)
    put(tf, t("문제는 정보의 부재가 아니라 역산 방법에 있다 → 순방향 연산만 수행하여 특이점이 "
              "원리적으로 발생하지 않고 O(1) 지연으로 동작하는 학습 기반 디커플러 채택",
              "The problem lies in the inversion method, not in missing information → adopt "
              "a learned decoupler that performs only forward operations, so singularities "
              "cannot arise, and runs with O(1) latency"),
        t(14, 13), False, DARK, space_after=0, line=1.35)

    # ── 15. 접촉 영역 물리 ──────────────────────────────────────────────────
    s = d.slide(t("접촉 모달리티 추가 — 새로 나타난 물리 현상",
                  "Adding the Contact Modality — New Physics Appears"),
                t("히스테리시스 · 크립 · 교차 민감도",
                  "Hysteresis, creep and cross-sensitivity"))
    d.pic(s, "fig3_pressure.png", MARGIN, BODY_TOP + Cm(0.4), w=Cm(30.5))
    tf = textbox(s, MARGIN, BODY_TOP + Cm(6.4), FULLW, Cm(5))
    put(tf, t("깊이를 고정한 채 유지(dwell)하는 동안  ΔL 2.33 %p · ΔR 5.73 %p 추가 변화 "
              "→ 힘 제거 후 잔차는 각각 0.57 %p · 0.08 %p",
              "During a constant-depth dwell, ΔL drifts a further 2.33 %p and ΔR 5.73 %p "
              "→ residuals after unloading are only 0.57 %p and 0.08 %p"),
        14, False, DARK, first=True, space_after=8, bullet="•")
    put(tf, t("즉 대부분이 영구 손상이 아니라 가역적인 시간 지연 성분 → 표준선형고체(SLS) 모델과 일치",
              "So most of it is a reversible time-lag component, not permanent damage "
              "→ consistent with a standard linear solid (SLS) model"),
        14, False, DARK, space_after=8, bullet="•")
    put(tf, t("dwell 구간 지수 완화 피팅 →  τ = 0.93 s   "
              "(채널별: τ_L 1.33 s, τ_R 0.59 s, 통합 τ = 1.01 s)",
              "Exponential relaxation fit over the dwell →  τ = 0.93 s   "
              "(channel-wise: τ_L 1.33 s, τ_R 0.59 s; combined τ = 1.01 s)"),
        14, True, MAROON, space_after=0, bullet="▶")

    # ── 16. Q-factor ────────────────────────────────────────────────────────
    s = d.slide(t("압축 시 인덕턴스가 증가해 보이는 이유",
                  "Why Inductance Appears to Rise Under Compression"),
                t("Q-factor 저하에 의한 역산 겉보기 효과",
                  "An inversion artefact caused by the falling Q factor"))
    shp, tf = rect(s, MARGIN, BODY_TOP + Cm(0.4), FULLW, Cm(1.7), LIGHT)
    put(tf, t("Q  =  ω · (최대 저장 에너지 / 주기당 소산 에너지)  =  X / R_s  =  ω L / R_s",
              "Q  =  ω · (max. stored energy / energy dissipated per cycle)  "
              "=  X / R_s  =  ω L / R_s"),
        t(17, 15), True, DARK, first=True, align=PP_ALIGN.CENTER, space_after=0)
    steps = [
        (t("센서가 눌린다", "The sensor is pressed"),
         "L ↓ ,  R_DC ↑ ,  R_eddy ↑", LIGHT, DARK),
        (t("Q 가 급감한다", "Q falls sharply"), "Q = ω L / R_s", LIGHT, DARK),
        (t("실제 공진주파수가 감소", "The true resonant frequency drops"),
         "f = 1/(2π√(LC)) · √(1 − 1/Q²)", LIGHT, DARK),
        (t("근사식으로 역산", "But we invert with the approximation"),
         "f ≈ 1/(2π√(LC)) ,  ΔL/L₀ = f₀²/f² − 1", LIGHT, DARK),
        (t("L 이 증가한 것처럼 보고됨", "L is reported as having increased"),
         t("겉보기 효과 — 실제 인덕턴스 증가가 아님",
           "an artefact — not a genuine rise in inductance"), PINK, MAROON),
    ]
    y = BODY_TOP + Cm(2.5)
    for t1, t2, fill, col in steps:
        shp, tf = rect(s, MARGIN, y, Cm(12.0), Cm(1.5), fill)
        put(tf, t1, t(15, 14), True, col, first=True, space_after=0)
        tf = textbox(s, MARGIN + Cm(12.8), y + Cm(0.35), Cm(17), Cm(1.5))
        put(tf, t2, t(14.5, 13.5), False, GREY, first=True, space_after=0)
        y += Cm(1.78)
    shp, tf = rect(s, MARGIN, y + Cm(0.3), FULLW, Cm(2.1), PINK)
    put(tf, t("본 연구는 이 항을 별도로 보정하지 않았다. 학습 기반 디커플러는 (ΔL, ΔR)에서 목표 "
              "물리량으로 가는 매핑을 직접 학습하므로, Q 의존성이 입력에 일관되게 포함되어 있는 한 "
              "역산 보정 없이도 올바른 출력을 낸다.",
              "No correction was applied. The learned decoupler maps (ΔL, ΔR) directly "
              "onto the target quantities, so as long as the Q dependence is consistently "
              "present in its inputs it produces correct outputs without an explicit "
              "inversion correction — a practical advantage over analytical inversion."),
        13.5, False, DARK, first=True, space_after=0, line=1.35)

    # ── 17. 모델 1 ──────────────────────────────────────────────────────────
    s = d.slide(t("디커플링 모델 ① — 2단계 구조와 경량화",
                  "Decoupling Model ① — Two-Stage Structure"),
                t("물리적 비대칭을 그대로 구조로 옮긴다",
                  "The physical asymmetry, transferred directly into the architecture"))
    shp, tf = rect(s, MARGIN, BODY_TOP + Cm(0.3), Cm(15.5), Cm(2.6), LIGHT)
    put(tf, "Stage 1 :  ΔR  →  ε̂          Stage 2 :  (ΔL, ε̂)  →  d̂",
        16, True, DARK, first=True, align=PP_ALIGN.CENTER, space_after=6)
    put(tf, t("“저항은 변형률만의 함수”라는 물리 사실을 귀납 편향으로 부여",
              "Imposes the physical fact that resistance is a function of strain alone, "
              "as an inductive bias"),
        12.5, False, GREY, align=PP_ALIGN.CENTER, space_after=0, line=1.25)
    tf = textbox(s, MARGIN, BODY_TOP + Cm(3.3), Cm(15.5), Cm(8.2))
    put(tf, t("동일 예산(≈ 50 K 파라미터) 공정 비교",
              "Fair comparison at an equal budget (≈ 50 K parameters)"),
        14, True, MAROON, first=True, space_after=8)
    put(tf, t("2단계 구조   근접 MAE 1.783 mm  (50,306 par.)",
              "Two-stage      proximity MAE 1.783 mm  (50,306 par.)"),
        13.5, False, DARK, space_after=5, bullet="•")
    put(tf, t("종단 구조     근접 MAE 1.827 mm  (49,592 par.)",
              "End-to-end    proximity MAE 1.827 mm  (49,592 par.)"),
        13.5, False, GREY, space_after=14, bullet="•")
    put(tf, t("오차 구조의 비대칭", "Asymmetry in the error structure"),
        14, True, MAROON, space_after=8)
    put(tf, t("변형률 오차는 동분산 — 전 구간에서 균일",
              "Strain error is homoscedastic — uniform across the range"),
        13.5, False, DARK, space_after=5, bullet="•")
    put(tf, t("근접 오차는 이분산 — d > 10 mm에서 급격히 불안정",
              "Distance error is heteroscedastic — unstable beyond d = 10 mm"),
        13.5, False, DARK, space_after=5, bullet="•")
    put(tf, t("→ 이후 모든 평가에서 d ≤ 15 mm, d ≤ 10 mm 구간을 별도 보고",
              "→ d ≤ 15 mm and d ≤ 10 mm are reported separately hereafter"),
        13, True, MAROON, space_after=0)
    d.pic(s, "fig4_pareto.png", MARGIN + Cm(16.5), BODY_TOP + Cm(0.5), w=Cm(13.5))
    tf = textbox(s, MARGIN + Cm(16.5), SH - Cm(2.6), Cm(14), Cm(1.6))
    put(tf, t("978 파라미터 medium-deep 이 파레토 무릎 — 50,306 파라미터 모델과 차이 0.09 mm",
              "The 978-parameter medium-deep sits at the Pareto knee — within 0.09 mm of "
              "the 50,306-parameter model"),
        12.5, False, GREY, first=True, space_after=0, line=1.3)

    # ── 18. 모델 2 ──────────────────────────────────────────────────────────
    s = d.slide(t("디커플링 모델 ② — SLS-EMA + Gate-MoE",
                  "Decoupling Model ② — SLS-EMA + Gate-MoE"),
                t("접촉/비접촉은 출력 물리량 자체가 다르다",
                  "Contact and non-contact differ in the very quantities to be produced"))
    d.pic(s, "fig5_moe.png", MARGIN + Cm(1.0), BODY_TOP + Cm(0.3), w=Cm(28.5))
    tf = textbox(s, MARGIN, BODY_TOP + Cm(6.0), FULLW, Cm(5))
    put(tf, t("SLS-EMA:  5장에서 구한 τ = 1.01 s 를 인과적 지수이동평균 이력 피처로 환원",
              "SLS-EMA: the τ = 1.01 s obtained in Section 5 is reduced to a causal "
              "exponential-moving-average history feature"),
        14.5, True, MAROON, first=True, space_after=7)
    put(tf, t("EMAₙ = EMAₙ₋₁ + α(xₙ − EMAₙ₋₁),   α = 1 − exp(−Δt/τ)   →   곱셈 1회 + 상태 변수 1개, "
              "MCU 비용 사실상 0",
              "EMAₙ = EMAₙ₋₁ + α(xₙ − EMAₙ₋₁),   α = 1 − exp(−Δt/τ)   →   one "
              "multiplication and one state variable; MCU cost effectively zero"),
        13.5, False, DARK, space_after=5, bullet="•")
    put(tf, t("정답 라벨이 아니라 실시간 ΔL·ΔR 만으로 계산 → 학습 피처를 배포 보드에서 그대로 재현 가능",
              "Computed from the live ΔL and ΔR alone, not from labels → the training "
              "feature is reproducible exactly on the deployed board"),
        13.5, False, DARK, space_after=5, bullet="•")
    put(tf, t("이 피처 하나로 접촉력 RMSE  0.521 N → 0.253 N  (F/T 센서 자체 잡음 ±0.19 N)",
              "This single feature cuts the force RMSE from 0.521 N to 0.253 N "
              "(the F/T sensor's own noise is ±0.19 N)"),
        13.5, True, MAROON, space_after=0, bullet="▶")

    # ── 19. 절제 실험 ───────────────────────────────────────────────────────
    s = d.slide(t("절제 실험", "Ablation Study"),
                t("Gate-MoE 와 SLS-EMA 의 기여를 분리한다 (D4, 사이클 단위 held-out)",
                  "Separating the contributions of Gate-MoE and SLS-EMA "
                  "(D4, cycle-wise held-out)"))
    hdr = [t("모델", "Model"), t("게이트 정확도", "Gate acc."),
           t("변형률 RMSE", "Strain RMSE"), t("근접거리 RMSE", "Distance RMSE"),
           t("접촉력 RMSE", "Force RMSE"), t("파라미터", "Param.")]
    data = [
        ("A:  Gate-MoE + SLS-EMA", "1.000", "0.108 %p", "0.318 mm", "0.253 N", "1,429"),
        (t("B:  SLS-EMA만 (통합)", "B:  SLS-EMA only (unified)"), "—", "0.288 %p",
         "0.322 mm", "0.079 mm †", "1,506"),
        (t("C:  Gate-MoE만 (이력 없음)", "C:  Gate-MoE only (no history)"), "0.982",
         "0.194 %p", "0.410 mm", "0.521 N", "1,317"),
        ("D:  Gate-MoE + GRU", "1.000", "0.123 %p", "0.318 mm", "0.263 N", "1,061"),
    ]
    cw = [Cm(8.4), Cm(4.4), Cm(4.4), Cm(4.6), Cm(4.4), Cm(4.4)]
    y = table(s, hdr, data, cw, MARGIN, BODY_TOP + Cm(0.5),
              hdr_h=Cm(1.35), row_h=Cm(1.25), hdr_size=t(12.5, 11.5),
              row_size=13.5, left_cols=(0,), accent_row=0)
    tf = textbox(s, MARGIN, y + Cm(0.25), FULLW, Cm(5))
    put(tf, t("이력 피처(SLS-EMA)의 이득은 접촉력에 집중  —  C→A 에서 힘 RMSE 0.521 → 0.253 N "
              "(절반 이하), 변형률은 0.194 → 0.108 %p",
              "The benefit of SLS-EMA is concentrated in force  —  from C to A the force "
              "RMSE falls 0.521 → 0.253 N (less than half), strain only 0.194 → 0.108 %p"),
        13.5, False, DARK, first=True, space_after=7, bullet="•")
    put(tf, t("게이트–전문가 구조의 기여는 B→A 비교에서  —  통합 모델은 변형률 RMSE 0.288 %p 로 "
              "2.7배 나쁨 (접촉/비접촉의 인장 응답이 다르기 때문)",
              "The gate-expert structure shows up from B to A  —  the unified model is "
              "2.7× worse in strain RMSE (0.288 %p), because the strain response differs "
              "between regimes"),
        13.5, False, DARK, space_after=7, bullet="•")
    put(tf, t("GRU(D)는 정확도 동률이면서 압력 전문가가 306 파라미터로 최소 — 다만 배포는 A 선택 "
              "(X-CUBE-AI Dense 지원 성숙, 은닉 상태 관리 불필요)",
              "The GRU (D) ties on accuracy with only 306 parameters for the pressure "
              "expert — yet A was deployed (mature dense-layer support in X-CUBE-AI, no "
              "hidden state to carry)"),
        13.5, False, DARK, space_after=0, bullet="•")

    # ── 20. 임베디드 ────────────────────────────────────────────────────────
    s = d.slide(t("임베디드 구현", "Embedded Implementation"),
                t("sklearn → ONNX → stedgeai → STM32,  매 TDM 주기마다 온보드 추론",
                  "sklearn → ONNX → stedgeai → STM32, with on-board inference once per "
                  "TDM cycle"))
    hdr = [t("항목", "Item"), t("2단 INT8 (초기)", "Two-stage INT8 (first)"),
           "medium-deep", t("MoE (최종)", "MoE (final)")]
    data = [
        (t("파라미터 수", "Parameters"), "50,306", "978", "1,413"),
        (t("가중치 ROM", "Weight ROM"), "~52 KB", "3,912 B", "5.66 KB"),
        (t("활성화 SRAM", "Activation SRAM"), "15.3 KB", "288 B", "368 B"),
        (t("전체 플래시", "Total flash"), "112.5 KB (85.8 %)", "~68 KB",
         "83.5 KB (65.2 %)"),
        (t("추론 지연 (DWT 실측)", "Inference latency (DWT)"), "1,068 µs", "141 µs",
         "260 µs (p95 263 µs)"),
        (t("1 ms 주기 내 완결", "Completes within 1 ms"), t("아니오", "no"),
         t("예", "yes"), t("예", "yes")),
    ]
    cw = [Cm(8.0), Cm(7.6), Cm(6.4), Cm(8.6)]
    y = table(s, hdr, data, cw, MARGIN, BODY_TOP + Cm(0.5),
              hdr_h=Cm(1.0), row_h=Cm(1.2), hdr_size=13, row_size=13.5,
              left_cols=(0,), accent_col=3)
    shp, tf = rect(s, MARGIN, y + Cm(0.25), FULLW, Cm(2.3), PINK)
    put(tf, t("지연 산포가 2 µs 수준  →  결정론적 제어 루프에 사용 가능",
              "The latency spread is only ~2 µs  →  usable in a deterministic control loop"),
        15, True, MAROON, first=True, space_after=6)
    put(tf, t("동일 추론을 PC에서 수행하면 USB 왕복 20–40 ms, 지터 ±5–15 ms  →  이 시스템에서 "
              "엣지 추론은 성능 최적화가 아니라 실시간 동작을 위한 요건",
              "The same inference on a PC costs a 20–40 ms USB round trip with ±5–15 ms "
              "jitter  →  edge inference here is a requirement for real-time operation, "
              "not an optimization"),
        13.5, False, DARK, space_after=0, line=1.3)

    # ── 21. 실시간 검증 ─────────────────────────────────────────────────────
    s = d.slide(t("실기 실시간 검증", "On-Hardware Real-Time Validation"),
                t("학습과 무관한 별도 세션 · PC 후처리 추론 없음 (세션 0819-153948, 809 s)",
                  "A session separate from training · no post-hoc inference on the PC "
                  "(session 0819-153948, 809 s)"))
    d.pic(s, "fig6_realtime.png", Cm((33.867 - 18.0) / 2), BODY_TOP + Cm(0.1),
          h=Cm(12.9))
    tf = textbox(s, MARGIN, SH - Cm(2.3), FULLW, Cm(1.8))
    put(tf, t("Part 1 (0–398 s) 변형률 고정 · 근접 반복 왕복    |    "
              "Part 2 (400–730 s) 근접 고정 · 변형률 스윕    |    Part 3 (738 s~) 두 축 동시",
              "Part 1 (0–398 s) strain held, proximity swept    |    "
              "Part 2 (400–730 s) proximity held, strain swept    |    "
              "Part 3 (from 738 s) both axes together"),
        13.5, False, GREY, first=True, align=PP_ALIGN.CENTER, space_after=0)

    # ── 22. 구간별 결과 ─────────────────────────────────────────────────────
    s = d.slide(t("구간별 결과", "Part-Wise Results"),
                t("각 축이 단독으로 변할 때와 동시에 변할 때를 분리해 관찰",
                  "Isolating the case where each axis varies alone from where both vary"))
    d.pic(s, "fig7_parity.png", MARGIN, BODY_TOP + Cm(0.4), w=Cm(16.5))
    hdr = [t("구간", "Segment"), t("변형률 R²", "Strain R²"),
           t("근접거리 R²", "Distance R²"), t("접촉력 R²", "Force R²")]
    data = [
        (t("Part 1 — 근접만 변화", "Part 1 — proximity only"), "0.981", "0.940", "0.430"),
        (t("Part 2 — 인장만 변화", "Part 2 — strain only"), "0.492", "−0.797", "—"),
        (t("Part 3 — 동시 변화", "Part 3 — simultaneous"), "0.602", "−2.051", "—"),
    ]
    cw = [Cm(6.0), Cm(2.9), Cm(3.1), Cm(2.6)]
    x0 = MARGIN + Cm(17.4)
    y = table(s, hdr, data, cw, x0, BODY_TOP + Cm(0.8),
              hdr_h=Cm(1.35), row_h=Cm(1.15), hdr_size=t(11.5, 11),
              row_size=12.5, left_cols=(0,), accent_row=0)
    tf = textbox(s, x0, y + Cm(0.3), Cm(14.6), Cm(8.2))
    put(tf, t("Part 1 — 세 물리량 모두 잘 복원",
              "Part 1 — all three quantities recovered well"),
        14, True, MAROON, first=True, space_after=6)
    put(tf, t("변형률 RMSE 1.37 %p · 근접 RMSE 2.11 mm\n게이트 정확도 98.3 %",
              "Strain RMSE 1.37 %p · distance RMSE 2.11 mm\nGate accuracy 98.3 %"),
        12.5, False, DARK, space_after=14, line=1.35)
    put(tf, t("Part 2 · 3 — 근접 R² 가 음수",
              "Parts 2 and 3 — the distance R² goes negative"),
        14, True, MAROON, space_after=6)
    put(tf, t("근접이 준정적이라 정답의 분산 자체가 매우 작다.\n"
              "R² = 1 − RMSE²/Var(y) 이므로 분모가 작으면\n"
              "절대 오차가 작아도 R² 가 크게 나빠진다.\n"
              "(Part 2 RMSE 6.19 mm — Part 1의 3배 수준)",
              "Proximity is quasi-static, so the ground-truth\nvariance is very small. "
              "Since R² = 1 − RMSE²/Var(y),\na small denominator degrades R² even for a "
              "modest\nabsolute error. (Part 2 RMSE 6.19 mm — about 3× Part 1)"),
        12.5, False, DARK, space_after=10, line=1.35)
    put(tf, t("그럼에도 저하는 실재 — 근접이 준정적일 때 인덕턴스에 남는 정보가 잡음 수준으로 줄어든다.",
              "The degradation is nonetheless real: when proximity is quasi-static, the "
              "information left in the inductance falls to the noise level."),
        12.5, True, DARK, space_after=0, line=1.35)

    # ── 23. 실패 모드 ───────────────────────────────────────────────────────
    s = d.slide(t("실기에서만 드러난 실패 모드",
                  "A Failure Mode Visible Only on Hardware"),
                t("오프라인 지표만으로는 임베디드 성능을 담보할 수 없다",
                  "Offline metrics alone cannot guarantee embedded performance"))
    shp, tf = rect(s, MARGIN, BODY_TOP + Cm(0.4), FULLW, Cm(1.6), PINK)
    put(tf, t("오프라인 held-out 게이트 정확도 100.0 %  →  그런데 실기에서 "
              "“뗀 뒤에도 계속 접촉으로 판정”",
              "Offline held-out gate accuracy 100.0 %  →  yet on hardware the gate kept "
              "reporting contact after release"),
        t(16, 15), True, MAROON, first=True, align=PP_ALIGN.CENTER, space_after=0,
        line=1.2)
    tf = textbox(s, MARGIN, BODY_TOP + Cm(2.4), FULLW, Cm(1.2))
    put(tf, t("원인 — 학습 조건과 배포 조건의 불일치가 만든 자기강화 루프",
              "Cause — a self-reinforcing loop created by the training/deployment mismatch"),
        15, True, DARK, first=True, space_after=0)
    loop = [
        t("접촉을 해제하면 원시 ΔL·ΔR 은 즉시 변한다",
          "On release, the raw ΔL and ΔR change immediately"),
        t("그러나 EMA 는 τ ≈ 1 s 동안 “눌린 값”에 머문다",
          "But the EMA stays at the pressed value for about τ ≈ 1 s"),
        t("게이트가 그 오래된 EMA 를 함께 보고 계속 “접촉”으로 판단한다",
          "The gate, seeing that stale EMA, keeps reporting contact"),
        t("판단이 바뀌지 않으니 EMA 리셋 조건도 걸리지 않는다  ↺",
          "Its decision never changes, so the EMA reset never triggers  ↺"),
    ]
    y = BODY_TOP + Cm(3.6)
    for i, txt in enumerate(loop):
        col = MAROON if i == 3 else HEADBG
        shp, tf = rect(s, MARGIN + Cm(0.6), y, Cm(1.1), Cm(1.1), col)
        put(tf, str(i + 1), 13, True, WHITE, first=True, align=PP_ALIGN.CENTER,
            space_after=0)
        tf = textbox(s, MARGIN + Cm(2.2), y + Cm(0.22), Cm(29), Cm(1.1))
        put(tf, txt, t(14.5, 14), i == 3, DARK if i < 3 else MAROON, first=True,
            space_after=0)
        y += Cm(1.35)
    tf = textbox(s, MARGIN, y + Cm(0.25), FULLW, Cm(1.2))
    put(tf, t("학습 시엔 정답 라벨(z)로 EMA 를 정확한 순간에 리셋할 수 있었지만, 배포 보드엔 그 정답이 없다.",
              "In training the EMA could be reset at the exact instant using the label z; "
              "the deployed board has no such ground truth."),
        13.5, False, GREY, first=True, space_after=0)
    shp, tf = rect(s, MARGIN, y + Cm(1.6), FULLW, Cm(3.1), BLUEBG)
    put(tf, t("해결 — 게이트를 EMA 없는 2입력 구조로 되돌림",
              "Fix — revert the gate to a two-input form without EMA"),
        15, True, C_V, first=True, space_after=6)
    put(tf, t("접촉 여부는 본래 즉각적인 물리적 전환이라 평활할 이유가 없다. 오프라인 정확도는 "
              "100.0 % → 98.2 % 로 낮아졌지만 실기 동작은 정상화되었다.\n"
              "회귀 정확도에 실제로 도움이 되는 전문가 A·B 는 4입력 유지.",
              "Contact is an instantaneous physical transition with no reason to be "
              "smoothed. Offline accuracy dropped from 100.0 % to 98.2 %, but behaviour "
              "on hardware was restored.\nExperts A and B, where history genuinely helps "
              "regression, kept their four inputs."),
        13, False, DARK, space_after=0, line=1.3)

    # ── 24. 결론 ────────────────────────────────────────────────────────────
    s = d.slide(t("결론", "Conclusion"))
    items = [
        (t("단일 전극에서 3종 신호 동시 취득",
           "Three signals acquired simultaneously from one electrode"),
         t("두 가닥 EGaIn 나선 코일 · 1 ms TDM · 1 kHz · CPU 유휴율 93.9 %",
           "Two-wire EGaIn spiral coil · 1 ms TDM · 1 kHz · 93.9 % CPU idle")),
        (t("디커플링의 물리적 근거를 실측으로 확인",
           "The physical basis for decoupling confirmed by measurement"),
         t("R 은 ε 에만, L 은 (ε, d) 에 의존 — ∇L 과 ∇R 이 전역적으로 선형 독립",
           "R depends on ε alone, L on (ε, d) — ∇L and ∇R are globally linearly "
           "independent")),
        (t("접촉 영역 물리를 규명하고 모델에 환원",
           "Contact-regime physics identified and folded into the model"),
         t("SLS 점탄성 τ ≈ 1.0 s → 인과적 EMA 이력 피처 / Q-factor 겉보기 효과 규명",
           "SLS viscoelasticity τ ≈ 1.0 s → causal EMA history feature; the Q-factor "
           "artefact explained")),
        (t("1,413 파라미터 모델을 STM32 에 임베딩",
           "A 1,413-parameter model embedded on the STM32"),
         t("매 TDM 주기마다 온보드 추론 260 µs (p95 263 µs), 지연 산포 약 2 µs",
           "On-board inference at 260 µs per TDM cycle (p95 263 µs), ~2 µs spread")),
        (t("실기 성능을 한계와 함께 보고",
           "On-hardware performance reported together with its limits"),
         t("Part 1: 변형률 R² 0.981 · 근접 R² 0.940 · 게이트 98 % / 준정적 구간 저하는 그대로 보고",
           "Part 1: strain R² 0.981 · distance R² 0.940 · gate 98 %; quasi-static "
           "degradation reported as measured")),
    ]
    y = BODY_TOP + Cm(0.6)
    for t1, t2 in items:
        rect(s, MARGIN, y, Cm(0.35), Cm(1.9), MAROON, radius=False)
        tf = textbox(s, MARGIN + Cm(0.9), y + Cm(0.1), FULLW - Cm(1.2), Cm(1.9))
        put(tf, t1, t(17, 16), True, DARK, first=True, space_after=5)
        put(tf, t2, 13.5, False, GREY, space_after=0, line=1.3)
        y += Cm(2.25)

    # ── 25. 향후 계획 ───────────────────────────────────────────────────────
    s = d.slide(t("향후 계획", "Future Work"))
    plans = [
        ("(i)", t("AC 등가 직렬 저항 R_s 추가 취득",
                  "Also acquire the AC equivalent series resistance R_s"),
         t("같은 TDM 슬롯에서 R_s 까지 얻으면 응답 시그니처가 완성되어 도체/유전체를 부호로 구분할 수 "
           "있고, Q-factor 보정도 직접 가능해진다.",
           "Obtaining R_s in the same TDM slot completes the response signature, letting "
           "conductors and dielectrics be told apart by sign and enabling the Q-factor "
           "correction to be applied directly.")),
        ("(ii)", t("다중 주파수 임피던스 측정", "Multi-frequency impedance measurement"),
         t("여러 주파수에서 임피던스 실수부·허수부를 얻으면 재질 구분과 근접 추정을 동시에 개선. "
           "준정적 구간 문제도 추가 채널로 완화될 여지.",
           "Real and imaginary impedance at several frequencies would improve material "
           "discrimination and proximity estimation together, and the extra channels may "
           "ease the quasi-static problem.")),
        ("(iii)", t("2차원 변형 분해", "Two-dimensional strain decomposition"),
         t("ΔL = f(Δl, ΔA, Δd) 인 반면 ΔR = f(Δl) → L 과 R 이 서로 다른 유효 고유 방향을 갖는다. "
           "십자형 시편과 4축 대칭 스테이지는 이미 준비 완료.",
           "ΔL = f(Δl, ΔA, Δd) whereas ΔR = f(Δl) → L and R have different effective "
           "eigen-directions. The cruciform specimen and four-axis stage are already in "
           "place.")),
        ("(iv)", t("순환 구조로의 전환", "Migration to a recurrent structure"),
         t("GRU 가 306 파라미터로 EMA 기반과 동등한 정확도. X-CUBE-AI 순환 레이어 지원 · 은닉 상태 "
           "리셋 · 양자화 안정성을 보드에서 검증하는 것이 다음 과제.",
           "The GRU matched the EMA-based models with 306 parameters. Verifying "
           "X-CUBE-AI's recurrent-layer support, hidden-state reset and quantization "
           "stability on this board is the next task.")),
    ]
    y = BODY_TOP + Cm(0.6)
    for num, t1, t2 in plans:
        shp, tf = rect(s, MARGIN, y, Cm(2.1), Cm(2.5), MAROON)
        put(tf, num, 17, True, WHITE, first=True, align=PP_ALIGN.CENTER, space_after=0)
        tf = textbox(s, MARGIN + Cm(2.7), y + Cm(0.2), FULLW - Cm(3.0), Cm(2.5))
        put(tf, t1, 16, True, DARK, first=True, space_after=6)
        put(tf, t2, 13, False, GREY, space_after=0, line=1.35)
        y += Cm(2.9)

    # ── 26. 마무리 ──────────────────────────────────────────────────────────
    s = d.slide()
    rect(s, Cm(0), Cm(0), SW, Cm(0.45), MAROON, radius=False)
    tf = textbox(s, MARGIN, Cm(7.2), FULLW, Cm(3))
    put(tf, t("감사합니다", "Thank you"), 40, True, DARK, first=True,
        align=PP_ALIGN.CENTER, space_after=16)
    put(tf, t("문희준  ·  고려대학교 기계공학과  BioRobotics & Control Lab",
              "Heejun Moon  ·  BioRobotics and Control Lab,  Korea University"),
        15, False, GREY, align=PP_ALIGN.CENTER, space_after=0)

    fn = t("발표자료_국문_단일전극_멀티모달센서.pptx",
           "Slides_EN_SingleElectrode_MultimodalSensor.pptx")
    out = os.path.join(HERE, fn)
    d.prs.save(out)
    print("saved:", fn)
    return out


if __name__ == "__main__":
    for lang in ("ko", "en"):
        build(lang)
